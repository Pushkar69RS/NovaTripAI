"""Build the Phase-2 Review-1 deck from the mandated department template.

    uv run python scripts/build_deck.py

The template in docs/review1/template/ is the format: its slide order, its
placeholders, its fonts, its footer and its slide numbering are kept exactly as
they are. Slides are only ever reused or duplicated from it, never restyled.
Where one mandated section needs more than one slide, the extra slide is a
duplicate of the same template slide with "(contd.)" in the title.

Every figure on a slide comes from docs/review1/numbers.md (the JSON block at
its end), which records the command each number came from; the literature
figures are attributed to their papers on the slide itself.
"""

from __future__ import annotations

import copy
import json
import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = (
    ROOT
    / "docs/review1/template"
    / "ISE-Division 5 7th sem Major Project PPT Template_Phase 2 review-1.pptx"
)
OUT = ROOT / "docs/review1/Travel_Yantra_Phase2_Review1.pptx"
SHOTS = ROOT / "docs/screenshots"
NUMBERS = ROOT / "docs/review1/numbers.md"


def load_numbers() -> dict:
    """Every figure on a slide comes from the JSON block in numbers.md."""
    text = NUMBERS.read_text(encoding="utf-8")
    return json.loads(text.split("```json", 1)[1].split("```", 1)[0])


N = load_numbers()
DEMO, RET, DB, COLD, LIT = (
    N["demo"],
    N["retrieval"],
    N["db"],
    N["cold_start"],
    N["literature"],
)

#: The template's own theme colours (theme1.xml), used for every diagram so no
#: colour enters the deck that the template did not already define.
INK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD1, 0x28, 0x2E)  # dk2
SAND = RGBColor(0xC8, 0xC8, 0xB1)  # lt2
GREY = RGBColor(0x7A, 0x7A, 0x7A)  # accent1
YELLOW = RGBColor(0xF5, 0xC2, 0x01)  # accent2
BLUE = RGBColor(0x52, 0x6D, 0xB0)  # accent3
MUTED = RGBColor(0x98, 0x9A, 0xAC)  # accent4
ORANGE = RGBColor(0xDC, 0x59, 0x24)  # accent5
OLIVE = RGBColor(0xB4, 0xB3, 0x92)  # accent6
SERIF = "Times New Roman"

# Template slide indices, by the section they carry.
T_TITLE, T_CONTENT, T_ABSTRACT, T_INTRO, T_PROBLEM = 0, 1, 2, 3, 4
T_OBJECTIVES, T_METHOD, T_DEMO, T_RESULTS, T_STATUS, T_PUB = 5, 6, 7, 8, 9, 10


# --------------------------------------------------------------------------- #
# template surgery
# --------------------------------------------------------------------------- #


def duplicate(prs: Presentation, index: int):
    """A copy of a template slide, appended at the end.

    Only slides whose shapes are placeholders are duplicated (every slide but
    the title one), so there are no image relationships to re-link.
    """
    source = prs.slides[index]
    new = prs.slides.add_slide(source.slide_layout)
    for shape in list(new.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        # Only the placeholders: a diagram already drawn on the source slide
        # must not follow its copy.
        if shape.is_placeholder:
            new.shapes._spTree.append(copy.deepcopy(shape._element))
    return new


def reorder(prs: Presentation, order: list[int]) -> None:
    """Rewrite the slide order. `order` holds current indices, in the order wanted."""
    ids = prs.slides._sldIdLst
    entries = list(ids)
    for entry in entries:
        ids.remove(entry)
    for index in order:
        ids.append(entries[index])


def body_of(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            return shape
    return None


def title_of(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            return shape
    return None


def set_title(slide, text: str) -> None:
    shape = title_of(slide)
    frame = shape.text_frame
    first = frame.paragraphs[0]
    for para in list(frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    if first.runs:
        first.runs[0].text = text
        for run in list(first.runs)[1:]:
            run._r.getparent().remove(run._r)
    else:
        first.add_run().text = text


def set_body(slide, items, size: float | None = None, space: int | None = None):
    """Replace the body text, cloning the template's own first paragraph.

    `items` are (level, text) or (level, text, bold). Cloning keeps the
    template's bullet character, indent, font and colour; only the size is
    overridden, and only where the slide would otherwise overflow.
    """
    shape = body_of(slide)
    frame = shape.text_frame
    frame.word_wrap = True
    paragraphs = list(frame.paragraphs)
    prototype = copy.deepcopy(paragraphs[0]._p)
    for para in paragraphs:
        para._p.getparent().remove(para._p)
    for item in items:
        level, text = item[0], item[1]
        bold = item[2] if len(item) > 2 else None
        element = copy.deepcopy(prototype)
        frame._txBody.append(element)
        para = frame.paragraphs[-1]
        for run in list(para.runs)[1:]:
            run._r.getparent().remove(run._r)
        if not para.runs:
            para.add_run()
        para.runs[0].text = text
        para.level = level
        if size is not None:
            para.runs[0].font.size = Pt(size)
        if bold is not None:
            para.runs[0].font.bold = bold
        if space is not None:
            para.space_after = Pt(space)
    return shape


def head(slide, text, *, size=None, top=0.17, height=0.82, left=0.5, width=8.96):
    """Set a section title and its box together. The template's own title boxes
    are 1.5 in tall for one line of text; a shorter box needs a matching size,
    or the text overflows upward out of the slide."""
    set_title(slide, text)
    shape = title_of(slide)
    place(shape, left=left, top=top, width=width, height=height)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    if size is not None:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(size)
    return shape


def place(shape, left=None, top=None, width=None, height=None) -> None:
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)


# --------------------------------------------------------------------------- #
# diagram primitives, all in the template's own colours and typeface
# --------------------------------------------------------------------------- #


def box(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    fill=WHITE,
    line=GREY,
    color=INK,
    size=10.5,
    bold=False,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    align=PP_ALIGN.CENTER,
    line_width=1.0,
    anchor=MSO_ANCHOR.MIDDLE,
):
    node = slide.shapes.add_shape(
        shape, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = fill
    if line is None:
        node.line.fill.background()
    else:
        node.line.color.rgb = line
        node.line.width = Pt(line_width)
    node.shadow.inherit = False
    frame = node.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = Emu(36000)
    frame.margin_top = frame.margin_bottom = Emu(54000)
    for i, line_text in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line_text
        run.font.size = Pt(size if i == 0 else size - 1)
        run.font.bold = bold and i == 0
        run.font.name = SERIF
        run.font.color.rgb = color
    return node


def label(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=9,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    italic=False,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = tb.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for i, line_text in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line_text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = SERIF
        run.font.color.rgb = color
    return tb


def arrow(slide, x1, y1, x2, y2, color=GREY, width=1.25):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    tail = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn

    end = tail.makeelement(
        qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}
    )
    tail.append(end)
    return line


def bar(slide, left, top, width, height, *, fill, text, size=9, color=WHITE):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    rect.shadow.inherit = False
    frame = rect.text_frame
    frame.margin_left = frame.margin_right = Emu(36000)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = SERIF
    run.font.color.rgb = color
    return rect


def table(slide, left, top, width, rows, *, col_widths, size=10.5, head_size=10.5):
    shape = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(0.3)
    )
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = True
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(0.26)
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Emu(54000)
            cell.margin_right = Emu(36000)
            cell.margin_top = cell.margin_bottom = Emu(9000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(value)
            run.font.size = Pt(head_size if r == 0 else size)
            run.font.bold = r == 0
            run.font.name = SERIF
    return tbl


def shot(slide, name, left, top, width, caption, *, size=9.5):
    """One screenshot with the single line that says what it proves."""
    path = SHOTS / name
    if not path.exists():
        raise SystemExit(f"missing screenshot: {path}")
    picture = slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width)
    )
    picture.line.color.rgb = GREY
    picture.line.width = Pt(0.75)
    height = Emu(picture.height).inches
    label(
        slide,
        left,
        top + height + 0.04,
        width,
        0.38,
        caption,
        size=size,
        color=RGBColor(0x33, 0x33, 0x33),
    )
    return picture


# --------------------------------------------------------------------------- #
# the deck
# --------------------------------------------------------------------------- #

TITLE = (
    "Travel Yantra — AI Based Personalized Travel Itinerary Planner "
    "and Virtual Tour Guide with Multilingual Support"
)
TEAM = [
    "Nithin G (1BY23IS140)",
    "Pushkar Reddy S (1BY23IS167)",
    "Rishab Paul (1BY23IS175)",
    "Rohan Balu (1BY23IS177)",
]


def slide_title(prs) -> None:
    slide = prs.slides[T_TITLE]
    body = body_of(slide)
    paragraphs = body.text_frame.paragraphs
    paragraphs[3].runs[0].text = TITLE
    paragraphs[3].runs[0].font.size = Pt(24)
    for i, member in enumerate(TEAM):
        paragraphs[5 + i].runs[0].text = member
    paragraphs[14].runs[0].text = "Prof. Shwetha T"
    paragraphs[15].runs[0].text = "Assistant Professor, Dept. of ISE"


def slide_content(prs) -> None:
    """The template's own agenda, unchanged: it is the mandated running order."""


def slide_abstract(prs) -> None:
    slide = prs.slides[T_ABSTRACT]
    set_body(
        slide,
        [
            (
                0,
                (
                    "Context: Planning a trip in India means a dozen browser tabs, almost "
                    "all in English, and a chat answer that may be confidently wrong about "
                    "when a place closes. At the site itself there is no guide that can be "
                    "trusted on a date or a fee."
                ),
            ),
            (
                0,
                (
                    "Problem Identified: Generic large language models are fluent but "
                    "unreliable planners — GPT-4 passes 0.6 percent of multi-constraint "
                    "plans on the TravelPlanner benchmark — and free-form generation "
                    "invents opening hours, fees and history that a traveller cannot check."
                ),
            ),
            (
                0,
                (
                    "Objective: This review demonstrates two of the six project "
                    "objectives in working software: Objective 1, AI-Powered Personalised "
                    "Itinerary Generation, and Objective 5, Virtual Cultural Tour Guide, "
                    "delivered as the feature named Katha."
                ),
            ),
            (
                0,
                (
                    "Methods: Our own solver computes the itinerary — candidate selection, "
                    "grouping places by part of the city, nearest-next ordering then "
                    "untangling under opening hours, a checklist, and a repair loop that "
                    "rebuilds one day. The language model reads the traveller's words into "
                    "the form, drafts places for a city we have never seen (labelled "
                    "unverified), writes the finished plan up and reads the guide aloud; "
                    "every paragraph it speaks is checked against its source."
                ),
            ),
            (
                0,
                (
                    f"Results: On the demonstration trip the routed itinerary is "
                    f"{DEMO['routed_km']} km against {DEMO['listed_km']} km for the same "
                    f"stops in list order, with {DEMO['checks_passed']} of "
                    f"{DEMO['checks_total']} checks passed in {DEMO['build_ms']} ms. "
                    f"Retrieval reaches Recall@5 of {RET['written']['hybrid'][0]:.3f} on "
                    f"written questions and {RET['lookups']['hybrid'][0]:.3f} on name "
                    f"lookups; {N['narration']['passed']} of {N['narration']['total']} "
                    f"narrated segments passed the fact-check; {N['tests']} automated "
                    "tests pass."
                ),
            ),
        ],
        size=14,
        space=8,
    )


def slide_introduction(prs) -> None:
    slide = prs.slides[T_INTRO]
    head(slide, "INTRODUCTION", size=30)
    set_body(
        slide,
        [
            (
                0,
                (
                    "Travel Yantra is an AI-based personalised itinerary planner and "
                    "virtual tour guide for the Indian tourism context, built for "
                    "travellers who currently plan across several disconnected platforms "
                    "and who are more comfortable in an Indian language than in English."
                ),
            ),
            (
                0,
                (
                    "Phase 1 established the problem, surveyed twenty-two papers and set "
                    "six measurable objectives. Phase 2 is implementation. This review "
                    "presents working software for two of those six objectives, evaluated "
                    "on real data rather than described as intent."
                ),
            ),
            (
                0,
                (
                    "Scope of this review — Objective 1: the planner is an India-wide "
                    "architecture; nothing in the solver is specific to one state. It is "
                    "demonstrated on Karnataka data because that is where our verified "
                    "data exists; a city outside it is drafted by the model on first "
                    "request and labelled unverified until a source confirms it."
                ),
            ),
            (
                0,
                (
                    "Scope of this review — Objective 5: the Katha guide is scoped to "
                    "Karnataka, since every paragraph it can speak has to be written and "
                    "sourced before it can be retrieved."
                ),
            ),
            (
                0,
                (
                    "Related work in one line: TravelPlanner (Xie et al., ICML 2024) "
                    "showed that a single language model fails multi-constraint travel "
                    "planning; Hao et al. (NAACL 2025) recovered the accuracy on that "
                    "same benchmark by handing the constraints to a solver; and "
                    "ChinaTravel (Shao et al., ICLR 2026) reproduced the effect on a "
                    "second, harder benchmark. Our architecture follows that finding "
                    "rather than restating the problem."
                ),
            ),
        ],
        size=16,
        space=10,
    )
    place(body_of(slide), left=0.5, top=1.45, width=8.96, height=5.4)


def slide_problem(prs) -> None:
    slide = prs.slides[T_PROBLEM]
    head(slide, "PROBLEM STATEMENT", size=30)
    set_body(
        slide,
        [
            (
                0,
                (
                    "A traveller assembling a three-day trip must reconcile opening "
                    "hours, weekly closing days, entry fees, road time between stops, a "
                    "midday meal, a budget and how far the oldest member of the group can "
                    "walk. These are hard constraints and they interact: fixing one "
                    "breaks another."
                ),
            ),
            (
                0,
                (
                    "Booking and review platforms each solve one slice — flights, hotels, "
                    "reviews, maps — and leave the traveller to be the integration layer. "
                    "Not one of them sequences a day against opening hours."
                ),
            ),
            (
                0,
                (
                    "General-purpose chat models are fluent, but measured at a 0.6 percent "
                    "pass rate on multi-constraint travel plans, and they state fees "
                    "and timings with a confidence they have not earned."
                ),
            ),
            (
                0,
                (
                    "Both classes of tool are English-first, which excludes the traveller "
                    "who would rather read and listen in Kannada or Hindi."
                ),
            ),
            (
                0,
                (
                    "On site, there is no grounded, source-cited guide. This is exactly "
                    "where a fabricated date or fee is least detectable by the listener "
                    "and most damaging to trust."
                ),
            ),
            (
                0,
                (
                    "Stated as an engineering problem: an itinerary must be computed "
                    "against its constraints and be able to show the working — not "
                    "generated and hoped for."
                ),
                True,
            ),
        ],
        size=16,
        space=10,
    )
    place(body_of(slide), left=0.5, top=1.45, width=8.96, height=5.4)


def slide_objectives(prs) -> None:
    slide = prs.slides[T_OBJECTIVES]
    head(slide, "OBJECTIVES OF THE WORK", size=30)
    set_body(
        slide,
        [
            (0, "Demonstrated in this review", True),
            (
                1,
                (
                    "Objective 1 — AI-Powered Personalised Itinerary Generation: "
                    "“Design and implement a multi-agent AI system that generates "
                    "day-by-day itineraries from user preferences, budget, dates and "
                    "group composition, targeting a user satisfaction rate of ninety "
                    "percent or higher.”"
                ),
            ),
            (
                1,
                (
                    "Objective 5 — Virtual Cultural Tour Guide: “Build a "
                    "Retrieval-Augmented Generation based interactive guide for history, "
                    "food, cultural significance and emergency information, with "
                    "sub-three-second response latency.” Delivered as the feature named "
                    "Katha."
                ),
            ),
            (
                0,
                (
                    "Scope: Objective 1 is an India-wide architecture demonstrated on "
                    "Karnataka; Objective 5 is scoped to Karnataka."
                ),
                True,
            ),
            (0, "Deferred to Review 2", True),
            (
                1,
                (
                    "Objective 2 — Multilingual Natural Language Support across Hindi, "
                    "Tamil, Telugu, Kannada and Bengali. English and Kannada are working "
                    "today; the full five-language pipeline is Review-2 work."
                ),
            ),
            (
                1,
                (
                    "Objective 3 — Multi-Channel Deployment and Accessibility across web, "
                    "WhatsApp and voice. The web channel is built; WhatsApp and telephony "
                    "are Review-2 work."
                ),
            ),
            (
                1,
                (
                    "Objective 4 — Real-Time Booking and Route Optimisation against live "
                    "booking and routing APIs."
                ),
            ),
            (
                0,
                (
                    "Objective 6 — Privacy and Cost Constraints is a standing constraint, "
                    "not a milestone: the stack is open-source-first and the measured "
                    "model spend for the whole project to date is under one US dollar."
                ),
            ),
        ],
        size=13.5,
        space=6,
    )
    place(body_of(slide), left=0.5, top=1.20, width=8.96, height=5.7)


# --------------------------------------------------------------------------- #
# methodology: the central claim, drawn
# --------------------------------------------------------------------------- #

PIPELINE = [
    ("1 Your words", "the model reads them into\na form you check", True),
    ("2 Unknown city?", "the model drafts its places,\nlabelled unverified", True),
    ("3 Candidates", "SQL + your interests +\npairing edges (our code)", False),
    (
        "4 Days and route",
        "group by part of the city;\nnearest-next, then untangle",
        False,
    ),
    ("5 Checklist", "hours, closures, day end,\nmeal, travel, budget", False),
    ("6 Repair", "drop one flexible stop,\nrebuild that day only", False),
    ("7 Reasons", "one plain sentence per\nstop, from templates", False),
    ("8 Narrate", "writes the plan up, reads\nKatha aloud, fact-checked", True),
]


def slide_methodology(prs) -> None:
    slide = prs.slides[T_METHOD]
    head(slide, "METHODOLOGY", size=30)
    set_body(
        slide,
        [
            (
                0,
                (
                    "The plan is computed by our own solver. The AI reads, drafts data "
                    "when we have none, and narrates. It never decides the plan."
                ),
                True,
            ),
            (
                0,
                (
                    "Stages 3 to 7 are ordinary code following fixed rules with a fixed "
                    "random seed, so the same request always gives the same plan."
                ),
            ),
        ],
        size=14,
        space=4,
    )
    place(body_of(slide), left=0.5, top=1.00, width=8.96, height=1.05)

    top, w, h, gap = 2.30, 2.12, 1.05, 0.14
    for row in (0, 1):
        y = top + row * (h + 0.80)
        for col in range(4):
            index = row * 4 + col
            stage, detail, is_llm = PIPELINE[index]
            x = 0.55 + col * (w + gap)
            box(
                slide,
                x,
                y,
                w,
                h,
                f"{stage}\n{detail}",
                fill=WHITE,
                line=ORANGE if is_llm else BLUE,
                line_width=1.75,
                size=9.5,
                bold=True,
            )
            if col < 3:
                arrow(slide, x + w, y + h / 2, x + w + gap, y + h / 2)
        if row == 0:
            drop = y + h + 0.40
            right = 0.55 + 3 * (w + gap) + w / 2
            arrow(slide, right, y + h, right, drop, color=GREY)
            arrow(slide, right, drop, 0.55 + w / 2, drop, color=GREY)
            arrow(slide, 0.55 + w / 2, drop, 0.55 + w / 2, y + h + 0.80, color=GREY)

    legend_y = 5.42
    box(slide, 0.55, legend_y, 0.22, 0.16, "", fill=BLUE, line=None)
    label(
        slide,
        0.85,
        legend_y - 0.03,
        3.5,
        0.25,
        "Our code — fixed rules, seeded",
        size=10,
    )
    box(slide, 4.35, legend_y, 0.22, 0.16, "", fill=ORANGE, line=None)
    label(
        slide,
        4.65,
        legend_y - 0.03,
        4.2,
        0.25,
        "Language model — reads, drafts, narrates",
        size=10,
    )
    label(
        slide,
        0.55,
        5.78,
        8.9,
        0.9,
        "The model also turns a chat edit into a strict instruction, and one day is "
        "rebuilt while the others stay untouched. Anything it drafts is labelled on "
        "screen until a source confirms it; anything it narrates is checked against "
        "its source, and a message it cannot read becomes one question, not a guess.",
        size=11,
        italic=True,
        color=RGBColor(0x33, 0x33, 0x33),
    )


def slide_methodology_2(slide) -> None:
    head(slide, "METHODOLOGY (contd.)", size=30)
    set_body(
        slide,
        [
            (
                0,
                (
                    "The validator is the contract. A plan is only returned once every "
                    "check passes, or the request is refused with the arithmetic shown."
                ),
                True,
            )
        ],
        size=14,
        space=2,
    )
    place(body_of(slide), left=0.5, top=1.00, width=8.96, height=0.5)

    # left: the repair flowchart
    label(
        slide,
        0.55,
        1.60,
        4.6,
        0.25,
        "Repair loop (maximum eight iterations)",
        size=11,
        bold=True,
    )
    cx, w, h = 0.55, 2.60, 0.5
    steps = [
        ("Build the day", BLUE),
        ("Validate: 6 checks per day", BLUE),
        ("Any violation?", GREY),
    ]
    y = 1.95
    for text, colour in steps:
        shape = MSO_SHAPE.DIAMOND if text.endswith("?") else MSO_SHAPE.ROUNDED_RECTANGLE
        box(
            slide,
            cx,
            y,
            w,
            h if shape != MSO_SHAPE.DIAMOND else 0.62,
            text,
            line=colour,
            line_width=1.5,
            size=10,
            shape=shape,
            bold=True,
        )
        y += 0.86
    box(
        slide,
        cx,
        y + 0.06,
        w,
        0.62,
        "Drop the lowest-scoring flexible\nstop of the least-repaired day",
        line=RED,
        line_width=1.5,
        size=9.5,
        bold=True,
    )
    box(
        slide,
        cx + w + 0.35,
        3.67,
        1.45,
        0.5,
        "Plan returned",
        line=BLUE,
        fill=SAND,
        line_width=1.5,
        size=10,
        bold=True,
    )
    arrow(slide, cx + w / 2, 2.45, cx + w / 2, 2.81)
    arrow(slide, cx + w / 2, 3.31, cx + w / 2, 3.67)
    arrow(slide, cx + w, 3.98, cx + w + 0.35, 3.92)
    label(slide, cx + w + 0.02, 3.66, 0.35, 0.2, "no", size=9)
    arrow(slide, cx + w / 2, 4.29, cx + w / 2, 4.61)
    label(slide, cx + w / 2 + 0.06, 4.33, 0.4, 0.2, "yes", size=9)
    arrow(slide, cx, 4.90, cx - 0.30, 4.90, color=RED)
    arrow(slide, cx - 0.30, 4.90, cx - 0.30, 2.20, color=RED)
    arrow(slide, cx - 0.30, 2.20, cx, 2.20, color=RED)
    label(
        slide,
        0.55,
        5.30,
        4.5,
        0.75,
        "The fixed-time anchor of a day is never dropped. Only the day named by the "
        "violation is rebuilt; every other day is left byte-for-byte identical.",
        size=10,
        italic=True,
        color=RGBColor(0x33, 0x33, 0x33),
    )

    # right: why a solver. Two benchmarks, two papers, one finding.
    label(
        slide,
        5.35,
        1.58,
        4.1,
        0.25,
        "Why a solver and not the model alone",
        size=11,
        bold=True,
    )
    # 100% of a benchmark maps to `span` inches. 0.6% would be invisible, so
    # every bar keeps a minimum sliver and its value is printed beside it.
    span = 2.45
    neural = LIT["chinatravel_neurosymbolic"] / LIT["chinatravel_neural_ratio"]
    groups = [
        (
            1.92,
            "TravelPlanner (Xie et al., ICML 2024) — final pass rate [1]",
            [
                ("GPT-4-Turbo", LIT["travelplanner_gpt4"], RED, None),
                ("LLM + solver [2]", LIT["travelplanner_solver"], BLUE, None),
            ],
        ),
        (
            2.98,
            "ChinaTravel (Shao et al., ICLR 2026) — constraint satisfaction [3]",
            [
                ("Purely neural", neural, RED, "about a tenth"),
                ("Neuro-symbolic", LIT["chinatravel_neurosymbolic"], BLUE, None),
            ],
        ),
    ]
    for top, caption, entries in groups:
        label(slide, 5.35, top, 4.1, 0.24, caption, size=8.5, italic=True, color=GREY)
        for i, (name, value, colour, shown) in enumerate(entries):
            y = top + 0.26 + i * 0.35
            width = max(value / 100 * span, 0.05)
            bar(slide, 6.72, y, width, 0.26, fill=colour, text="")
            label(slide, 5.35, y + 0.02, 1.30, 0.24, name, size=9, align=PP_ALIGN.RIGHT)
            label(
                slide,
                6.79 + width,
                y + 0.01,
                1.4,
                0.24,
                shown or f"{value:g}%",
                size=9.5,
                bold=True,
            )
    label(
        slide,
        5.35,
        4.02,
        4.1,
        1.5,
        f"Pure language-model agents score in the low single digits on both "
        f"benchmarks. Pairing the model with a solver raises that to "
        f"{LIT['travelplanner_solver']:g} percent on TravelPlanner (Hao et al.) and "
        f"tenfold on ChinaTravel, the papers' own comparison. These are published "
        f"figures, not ours, and they are why the plan here is computed by a solver "
        f"and the model is kept out of the decision.",
        size=10,
        color=INK,
    )
    label(
        slide,
        5.35,
        5.50,
        4.1,
        0.6,
        f"Our own result on the demonstration trip — {DEMO['checks_passed']} of "
        f"{DEMO['checks_total']} checks passed — is on the results slide.",
        size=10,
        bold=True,
    )


def slide_methodology_3(slide) -> None:
    head(slide, "METHODOLOGY (contd.)", size=30)
    set_body(
        slide,
        [
            (
                0,
                (
                    "Katha, the guide: the model speaks only from paragraphs we have "
                    "written and stored, and every paragraph carries its source."
                ),
                True,
            )
        ],
        size=14,
        space=2,
    )
    place(body_of(slide), left=0.5, top=1.00, width=8.96, height=0.5)

    label(
        slide,
        0.55,
        1.60,
        4.3,
        0.25,
        "City Katha — a fixed portrait",
        size=11,
        bold=True,
    )
    box(
        slide,
        0.55,
        1.90,
        4.3,
        3.15,
        "What the city is · how it began · who ruled here · the city today · "
        "what it eats · festivals · worth your time · before you come\n\n"
        "Each paragraph sits at a tier — 2, 5 or 10 minutes — so a longer Katha "
        "is a superset of a shorter one. More minutes means more of the city, "
        "never more of the inside of one monument; that lives in the Place Katha "
        "behind “Go deeper”.\n\n"
        "No retrieval and no randomness on this path: the same city and the same "
        "minutes always tell the same story, in the same order. A city we have "
        "never seen gets its portrait drafted by the model and labelled.",
        fill=WHITE,
        line=BLUE,
        line_width=1.5,
        size=10,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        shape=MSO_SHAPE.RECTANGLE,
    )

    label(
        slide,
        5.15,
        1.60,
        4.3,
        0.25,
        "Place and Day Kathas, and questions — retrieval",
        size=11,
        bold=True,
    )
    box(
        slide,
        5.15,
        1.90,
        4.3,
        3.15,
        "Every paragraph is turned into a list of numbers that captures its "
        "meaning, so a question asked in Kannada lands near an answer written in "
        "English. A second search matches exact words, for names. The two rankings "
        "are merged; below 0.81 similarity with no exact match the guide refuses "
        "rather than answers.\n\n"
        "After the model writes a paragraph, every date, number and name is checked "
        "against the source it was given. Anything of its own is thrown away and the "
        "model tries once more; if that fails too, the source text is spoken as it "
        "is.\n\n"
        "Speech is Sarvam, served from our own cache for the demo.",
        fill=WHITE,
        line=ORANGE,
        line_width=1.5,
        size=10,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        shape=MSO_SHAPE.RECTANGLE,
    )
    label(
        slide,
        0.55,
        5.25,
        8.9,
        0.5,
        "For the record: pgvector HNSW for the meaning search, tsvector for the exact "
        "words, reciprocal rank fusion (k = 60), multilingual-e5-small (384 "
        "dimensions), google/gemini-3.1-flash-lite for narration.",
        size=9,
        italic=True,
        color=GREY,
    )
    label(
        slide,
        0.55,
        5.80,
        8.9,
        0.7,
        "Rhythm rules for a Place or Day Katha: a paragraph never repeats; it opens "
        "on a hook or a story when there is one; five minutes or more carries at "
        "least one story.",
        size=10,
        italic=True,
        color=RGBColor(0x33, 0x33, 0x33),
    )


# --------------------------------------------------------------------------- #
# demonstration
# --------------------------------------------------------------------------- #

SHOT_ROWS = [
    [
        (
            "04b-form-step1.png",
            "The form, filled by the model from one pasted sentence; the traveller checks it.",
        ),
        (
            "07b-building.png",
            (
                f"Build trace from PlanMetrics: {DEMO['candidates']} of {DEMO['places']} candidates, "
                f"{DEMO['listed_km']} → {DEMO['routed_km']} km, {DEMO['checks_passed']}/{DEMO['checks_total']} checks."
            ),
        ),
        (
            "09b-plan-routed.png",
            "Plan page: “In a few words”, written by the model and checked; the getting-around line.",
        ),
        (
            "10b-plan-day2-listed.png",
            (
                f"Day 2 as listed, on real tiles: {DEMO['day2_listed_km']} km in list order against "
                f"{DEMO['day2_routed_km']} km routed."
            ),
        ),
    ],
    [
        (
            "12b-plan-chat-edit.png",
            "A chat edit rebuilds Day 2 only; Days 1 and 3 untouched, and the reply says so.",
        ),
        (
            "17b-katha-city.png",
            "Mysuru city Katha at 2 minutes: a fixed portrait, no type labels, a source per paragraph.",
        ),
        (
            "14b-katha-place.png",
            "Place Katha for Mysore Palace, with the depth picker; monument interiors live here.",
        ),
        (
            "15b-katha-home.png",
            "Katha home: one pin per region with its place count, on real tiles.",
        ),
    ],
    [
        (
            "07c-building-coldstart.png",
            f"Mangalore, never seen before: {COLD['places']} places drafted in {COLD['seconds']:g} s, labelled unverified.",
        ),
        (
            "09c-plan-mangalore.png",
            "The Mangalore plan: our solver scheduled the draft; each stop says “AI-drafted · unverified”.",
        ),
        (
            "19c-doesnt-fit.png",
            "Mysuru and Hampi in one day: refused with the arithmetic and two buildable alternatives.",
        ),
    ],
]


def slide_demo(slide, rows, part: int, total: int) -> None:
    set_title(
        slide,
        "DEMONSTRATION OF PROJECT EXECUTION"
        + ("" if part == 1 else f" (contd. {part}/{total})"),
    )
    place(title_of(slide), left=0.5, top=0.17, width=9.26, height=0.72)
    for run in title_of(slide).text_frame.paragraphs[0].runs:
        run.font.size = Pt(24)
    body = body_of(slide)
    body._element.getparent().remove(body._element)
    w = 3.95
    for i, (name, caption) in enumerate(rows):
        left = 0.55 + (i % 2) * (w + 0.5)
        top = 1.10 + (i // 2) * 2.92
        shot(slide, name, left, top, w, caption)


# --------------------------------------------------------------------------- #
# results
# --------------------------------------------------------------------------- #


def slide_results(prs) -> None:
    slide = prs.slides[T_RESULTS]
    head(slide, "Experimental Setup, Results & Analysis", size=28, height=0.72)
    set_body(
        slide,
        [
            (
                0,
                (
                    "Setup: Python 3.12 and FastAPI against Supabase PostgreSQL 17 with "
                    "pgvector and pg_trgm; embeddings computed locally with "
                    "multilingual-e5-small; models through OpenRouter; speech from Sarvam "
                    "bulbul:v3. Planner and retrieval are measured by their own harnesses, "
                    "offline where they can be, so the numbers are reproducible."
                ),
                False,
            )
        ],
        size=12.5,
        space=2,
    )
    place(body_of(slide), left=0.5, top=0.92, width=8.96, height=0.85)

    label(
        slide,
        0.55,
        1.82,
        4.2,
        0.25,
        "Data the system was built on",
        size=11.5,
        bold=True,
    )
    table(
        slide,
        0.55,
        2.10,
        4.15,
        [
            ["Dataset", "Count"],
            ["Places in the table", str(DB["poi_total"])],
            ["Verified against a live source", str(DB["poi_verified"])],
            ["Of those, corrected by verification", str(DB["corrected_of_verified"])],
            ["Drafted by the model, unverified", str(DB["poi_ai"])],
            ["Live corpus paragraphs", str(DB["chunks_live"])],
            ["Of which the city portrait layer", str(DB["chunks_themed"])],
        ],
        col_widths=[3.05, 1.10],
        size=10,
        head_size=10,
    )
    label(
        slide,
        0.55,
        3.92,
        4.15,
        0.62,
        f"{DB['corrected_of_verified']} of the {DB['poi_verified']} verified entries "
        "were wrong before we checked them. That ratio is the argument for verifying "
        "rather than trusting, and for labelling what is not yet verified.",
        size=9.5,
        italic=True,
        color=RGBColor(0x33, 0x33, 0x33),
    )

    label(slide, 0.55, 4.62, 4.15, 0.25, "Module-wise status", size=11.5, bold=True)
    table(
        slide,
        0.55,
        4.90,
        4.15,
        [
            ["Module", "State"],
            ["Intake parser and cold start", "Working"],
            ["Planner (cluster, route, validate, repair)", "Working"],
            ["Retrieval; city, place and day Kathas", "Working"],
            ["Narration with fact-check", "Working"],
            ["Speech (Sarvam, cached)", "Working"],
            ["Web application, maps and chat", "Working"],
        ],
        col_widths=[3.15, 1.00],
        size=9.5,
        head_size=9.5,
    )

    label(
        slide,
        5.10,
        1.82,
        4.35,
        0.25,
        "Planner, on the seeded demonstration trip",
        size=11.5,
        bold=True,
    )
    # One scale for all four bars, so Day 2 cannot be read as though it were
    # the whole trip. Labels sit outside the short bars.
    scale = 2.55 / DEMO["listed_km"]
    rows = [
        (2.34, "Whole trip, as listed", DEMO["listed_km"], RED),
        (2.70, "Whole trip, as routed", DEMO["routed_km"], BLUE),
        (3.16, "Day 2, as listed", DEMO["day2_listed_km"], RED),
        (3.52, "Day 2, as routed", DEMO["day2_routed_km"], BLUE),
    ]
    for y, name, value, colour in rows:
        bar(slide, 6.45, y, value * scale, 0.28, fill=colour, text="")
        label(slide, 5.10, y + 0.03, 1.30, 0.25, name, size=9.5, align=PP_ALIGN.RIGHT)
        label(
            slide,
            6.50 + value * scale,
            y + 0.02,
            1.05,
            0.25,
            f"{value:g} km",
            size=9.5,
            bold=True,
        )
    label(
        slide,
        5.10,
        3.86,
        4.35,
        0.22,
        "Same scale for all four bars.",
        size=8.5,
        italic=True,
        color=GREY,
    )
    table(
        slide,
        5.10,
        4.18,
        4.35,
        [
            ["Planner measure", "Result"],
            [
                "Constraint checks passed",
                f"{DEMO['checks_passed']} / {DEMO['checks_total']}",
            ],
            ["Repairs applied", str(DEMO["repairs"])],
            ["Candidates considered", f"{DEMO['candidates']} of {DEMO['places']}"],
            ["Build time", f"{DEMO['build_ms']} ms"],
        ],
        col_widths=[3.15, 1.20],
        size=10,
        head_size=10,
    )
    label(
        slide,
        5.10,
        5.68,
        4.35,
        0.85,
        "“As listed” is the same stops visited in candidate order — what an "
        "itinerary that never routes produces. It is the honest baseline, because "
        "nearest-next was already optimal here.",
        size=10,
        italic=True,
        color=RGBColor(0x33, 0x33, 0x33),
    )


def slide_results_2(slide) -> None:
    head(slide, "Experimental Setup, Results & Analysis (contd.)", size=24, height=0.66)
    body = body_of(slide)
    body._element.getparent().remove(body._element)

    w, lk, kn = RET["written"], RET["lookups"], RET["kannada"]
    label(
        slide,
        0.55,
        1.00,
        5.6,
        0.25,
        f"Retrieval — Recall@5, re-measured today (eval_run {RET['eval_run']})",
        size=11.5,
        bold=True,
    )
    table(
        slide,
        0.55,
        1.30,
        5.55,
        [
            ["Method", "30 questions", "10 name lookups", "Kannada", "p50"],
            [
                "Dense (pgvector, e5)",
                f"{w['dense'][0]:.3f}",
                f"{lk['dense'][0]:.3f}",
                f"{kn['dense']:.3f}",
                f"{w['dense'][2]:.0f} ms",
            ],
            [
                "Lexical (tsvector)",
                f"{w['lexical'][0]:.3f}",
                f"{lk['lexical'][0]:.3f}",
                f"{kn['lexical']:.3f}",
                f"{w['lexical'][2]:.0f} ms",
            ],
            [
                "Hybrid + RRF",
                f"{w['hybrid'][0]:.3f}",
                f"{lk['hybrid'][0]:.3f}",
                f"{kn['hybrid']:.3f}",
                f"{w['hybrid'][2]:.0f} ms",
            ],
        ],
        col_widths=[1.75, 1.05, 1.25, 0.75, 0.75],
        size=9.5,
        head_size=9.5,
    )
    label(
        slide,
        0.55,
        2.60,
        5.55,
        1.45,
        "Reported honestly: hybrid retrieval tied dense rather than beating it, at "
        "nearly twice the latency. We keep it for one measured reason — on exact "
        f"names the word search is as good as the meaning search ({lk['lexical'][0]:.3f} "
        f"against {lk['dense'][0]:.3f}) at less than half the time, which is how a "
        "traveller actually looks for a monument — and the meaning search carries "
        f"Kannada ({kn['dense']:.3f} against {kn['lexical']:.3f}), where the English "
        "index has almost nothing to match.",
        size=10,
    )
    label(
        slide,
        0.55,
        4.05,
        5.55,
        1.25,
        "Also reported: Recall@5 fell from 0.933 to 0.867 when the hand-written "
        "corpus was merged in, and was not tuned back. Re-measured after today's city "
        f"portrait layer: Recall@5 unchanged, MRR {RET['previous_mrr']:.3f} → "
        f"{w['hybrid'][1]:.3f}, p50 {RET['previous_p50']} → {w['hybrid'][2]:.0f} ms. "
        "The numbers stand as measured.",
        size=10,
        bold=True,
    )

    label(
        slide, 6.35, 1.00, 3.1, 0.25, "Grounding, refusal, tests", size=11.5, bold=True
    )
    table(
        slide,
        6.35,
        1.30,
        3.10,
        [
            ["Measure", "Result"],
            [
                "Narrated segments passing the check",
                f"{N['narration']['passed']} / {N['narration']['total']}",
            ],
            ["Refusal gate", f"{RET['refusal'][0]} / {RET['refusal'][1]}"],
            ["Automated tests", str(N["tests"])],
        ],
        col_widths=[2.05, 1.05],
        size=10,
        head_size=10,
    )
    label(
        slide,
        6.35,
        2.55,
        3.1,
        0.25,
        "Cold start — a city never seen",
        size=11.5,
        bold=True,
    )
    table(
        slide,
        6.35,
        2.83,
        3.10,
        [
            ["City", "Places", "Paragraphs", "Time", "Cost"],
            [
                COLD["city"],
                str(COLD["places"]),
                str(COLD["paragraphs"]),
                f"{COLD['seconds']:g} s",
                f"${COLD['cost_usd']:.3f}",
            ],
        ],
        col_widths=[0.90, 0.55, 0.75, 0.45, 0.45],
        size=9,
        head_size=9,
    )
    label(
        slide,
        6.35,
        3.55,
        3.10,
        1.75,
        "Every narrated segment is checked after generation against the paragraphs "
        "it was given: any date, number or English name not in the source fails, the "
        "segment is retried, then replaced by the source text itself.\n\n"
        "The refusal gate is the opposite test — six questions the corpus cannot "
        "answer, all six refused.",
        size=9.5,
    )
    label(
        slide,
        0.55,
        5.45,
        8.9,
        0.8,
        f"Analysis: the planner result is a checklist result, not an opinion — "
        f"{DEMO['checks_passed']} of {DEMO['checks_total']} checks on a trip with an "
        "elder, a child, a fixed-time palace slot and a 7 pm evening. The retrieval "
        "result is reported as measured, including the regression, because a "
        "benchmark that only moves upward is not a benchmark.",
        size=11,
        italic=True,
    )


def slide_status(prs) -> None:
    slide = prs.slides[T_STATUS]
    head(slide, "PROJECT STATUS, CONCLUSION", size=28, height=0.72)
    body = body_of(slide)
    body._element.getparent().remove(body._element)

    columns = [
        (
            "Completed",
            BLUE,
            (
                "Intake: your words read into a form\nCold start for an unseen city\n"
                "Deterministic planner\nHybrid retrieval\nCity, Place and Day Kathas\n"
                "Narration with fact-check\nSarvam speech, cached\nChat repair, one day\n"
                "Web app on real map tiles"
            ),
        ),
        (
            "In progress",
            YELLOW,
            (
                "Verifying the model-drafted places\n\nCounting getting-around cost\n"
                "against the budget\n\nOpenStreetMap as the place source\nfor new cities"
            ),
        ),
        (
            "Not started — Review 2",
            GREY,
            (
                "Objective 2: full five-\nlanguage support\n\nObjective 3: WhatsApp and\nvoice channels\n\n"
                "Objective 4: live booking\nand routing APIs\n\nA dedicated vector database\nfor the five-language corpus"
            ),
        ),
    ]
    for i, (column, colour, items) in enumerate(columns):
        x = 0.55 + i * 3.03
        box(
            slide,
            x,
            1.00,
            2.85,
            0.40,
            column,
            fill=colour,
            color=WHITE if colour is not YELLOW else INK,
            line=None,
            size=11.5,
            bold=True,
            shape=MSO_SHAPE.RECTANGLE,
        )
        box(
            slide,
            x,
            1.40,
            2.85,
            2.30,
            items,
            fill=WHITE,
            line=colour,
            size=10,
            align=PP_ALIGN.LEFT,
            shape=MSO_SHAPE.RECTANGLE,
            anchor=MSO_ANCHOR.TOP,
        )

    label(
        slide,
        0.55,
        3.85,
        8.9,
        0.25,
        "Known limits, stated plainly",
        size=11.5,
        bold=True,
    )
    limits = [
        (
            f"{DB['poi_ai']} places drafted by the model for {COLD['city']} are unverified; "
            "every one is labelled “AI-drafted · unverified” on screen until a source confirms it."
        ),
        (
            "The getting-around cost is shown per day, estimated, and not yet counted "
            "against the budget or checked by the planner."
        ),
        (
            f"{DB['poi_draft']} of the seeded places are still unverified draft data; the "
            "interface labels them as estimated rather than hiding it."
        ),
        (
            "Routes are drawn as straight segments between stops, on real tiles or on the "
            "sketch; the kilometres use a detour factor and are honest, the drawing is not."
        ),
        (
            "The system runs on localhost for this review; deployment is a post-review "
            "decision recorded in the project log."
        ),
    ]
    y = 4.12
    for text in limits:
        box(
            slide,
            0.55,
            y,
            0.10,
            0.10,
            "",
            fill=RED,
            line=None,
            shape=MSO_SHAPE.RECTANGLE,
        )
        label(slide, 0.78, y - 0.06, 8.65, 0.42, text, size=10)
        y += 0.40

    label(
        slide,
        0.55,
        6.15,
        8.9,
        0.7,
        "Conclusion: two of six objectives are demonstrated in working software, "
        "measured rather than asserted. Beyond Karnataka, a city we have never seen is "
        "drafted by the model on first request, labelled, and still scheduled by our "
        "solver; OpenStreetMap is the intended place source for Review 2.",
        size=11,
        bold=True,
    )


def slide_publication(prs) -> None:
    slide = prs.slides[T_PUB]
    head(slide, "PUBLICATION", size=30)
    set_body(
        slide,
        [
            (
                0,
                (
                    "Planned contribution: IndiaTravel — the first India-specific "
                    "benchmark for travel-planning AI agents."
                ),
                True,
            ),
            (
                1,
                (
                    "Modelled on ChinaTravel (Shao et al., ICLR 2026), the equivalent "
                    "open-ended benchmark for Chinese travel, on which neuro-symbolic "
                    "agents reach 37.0 percent constraint satisfaction against 2.6 "
                    "percent for purely neural ones — roughly a tenfold gap."
                ),
            ),
            (
                1,
                (
                    "The novelty is the constraint set. India-specific constraints have "
                    "no counterpart in the existing benchmarks: temple darshan windows and "
                    "midday closures, weekly closing days, festival and Dasara crowding, "
                    "monsoon access to hill roads, vegetarian and Jain food requirements, "
                    "multi-generational parties with an elder's walking limit, and "
                    "state-boundary transport changes."
                ),
            ),
            (
                1,
                (
                    f"Our planner, its validator and the {DB['poi_total']}-place dataset "
                    f"({DB['poi_verified']} verified against live sources) are the working "
                    "reference implementation the benchmark would be released alongside."
                ),
            ),
            (
                0,
                (
                    "Target: a Scopus-indexed conference or journal. Drafting begins "
                    "after Review 2, once the multilingual and multi-channel objectives "
                    "provide the second half of the evaluation."
                ),
                True,
            ),
            (
                0,
                (
                    "The report will be prepared in LaTeX and checked through Turnitin, "
                    "to remain under the fifteen percent similarity requirement."
                ),
            ),
        ],
        size=13,
        space=7,
    )
    place(body_of(slide), left=0.5, top=1.05, width=8.96, height=5.8)


REFERENCES = [
    (
        "J. Xie, K. Zhang, J. Chen, T. Zhu, R. Lou, Y. Tian, Y. Xiao and Y. Su, "
        "“TravelPlanner: A Benchmark for Real-World Planning with Language Agents,” "
        "in Proc. 41st Int. Conf. Machine Learning (ICML), 2024."
    ),
    (
        "Y. Hao, Y. Chen, Y. Zhang and C. Fan, “Large Language Models Can Solve "
        "Real-World Planning Rigorously with Formal Verification Tools,” in Proc. "
        "2025 Conf. North American Chapter of the ACL (NAACL), 2025, pp. 3499–3557."
    ),
    (
        "J.-J. Shao, B.-W. Zhang, X.-W. Yang, B. Chen, L.-Z. Guo and Y.-F. Li, "
        "“ChinaTravel: An Open-Ended Travel Planning Benchmark with Compositional "
        "Constraint Validation for Language Agents,” in Proc. Int. Conf. Learning "
        "Representations (ICLR), 2026."
    ),
    (
        "P. Lewis et al., “Retrieval-Augmented Generation for Knowledge-Intensive "
        "NLP Tasks,” in Adv. Neural Inf. Process. Syst. (NeurIPS), vol. 33, 2020, "
        "pp. 9459–9474."
    ),
    (
        "L. Wang, N. Yang, X. Huang, L. Yang, R. Majumder and F. Wei, “Multilingual "
        "E5 Text Embeddings: A Technical Report,” arXiv:2402.05672, 2024."
    ),
    (
        "G. V. Cormack, C. L. A. Clarke and S. Buettcher, “Reciprocal Rank Fusion "
        "Outperforms Condorcet and Individual Rank Learning Methods,” in Proc. 32nd "
        "Int. ACM SIGIR Conf., 2009, pp. 758–759."
    ),
    (
        "G. A. Croes, “A Method for Solving Traveling-Salesman Problems,” Operations "
        "Research, vol. 6, no. 6, pp. 791–812, 1958."
    ),
    (
        "D. Arthur and S. Vassilvitskii, “k-means++: The Advantages of Careful "
        "Seeding,” in Proc. 18th ACM-SIAM Symp. Discrete Algorithms (SODA), 2007, "
        "pp. 1027–1035."
    ),
    (
        "S. P. Lloyd, “Least Squares Quantization in PCM,” IEEE Trans. Information "
        "Theory, vol. 28, no. 2, pp. 129–137, 1982."
    ),
]


def slide_references(slide) -> None:
    head(slide, "REFERENCES (IEEE FORMAT)", size=28, height=0.70)
    set_body(
        slide,
        [(0, f"[{i}] {text}", False) for i, text in enumerate(REFERENCES, 1)],
        size=10,
        space=4,
    )
    place(body_of(slide), left=0.5, top=1.02, width=8.96, height=5.95)


#: Text that means a Python value leaked into a slide instead of its label.
LEAKED = ("<function", "0x0", "object at", "None")


def check(path: Path) -> None:
    """Fail loudly if any shape carries a repr rather than the text intended,
    and report any text box whose estimated height overruns its frame."""
    from pptx import Presentation as _Presentation

    prs = _Presentation(str(path))
    bad = [
        (i, shape.name, needle, shape.text_frame.text[:60])
        for i, slide in enumerate(prs.slides, 1)
        for shape in slide.shapes
        if shape.has_text_frame
        for needle in LEAKED
        if needle in shape.text_frame.text
    ]
    if bad:
        for entry in bad:
            print(f"  LEAKED {entry}")
        raise SystemExit(f"{len(bad)} shape(s) carry a leaked Python value")

    # Overflow estimate: characters per line from the frame width and the font
    # size (Times, about 0.5 em per character), lines at 1.2 em. Crude, so it
    # only has to be right about the ones that clearly do not fit.
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            width_in = shape.width / 914400
            height_in = shape.height / 914400
            lines = 0.0
            for para in shape.text_frame.paragraphs:
                size = next(
                    (r.font.size.pt for r in para.runs if r.font.size is not None), 18.0
                )
                cpl = max(8.0, width_in * 72 / (size * 0.5))
                text = "".join(r.text for r in para.runs)
                lines += max(1.0, math.ceil(len(text) / cpl)) * size * 1.2 / 72
            if lines > height_in * 1.25 + 0.1:
                print(
                    f"  OVERFLOW? slide {i} {shape.name!r}: ~{lines:.2f} in of text in "
                    f"{height_in:.2f} in — {shape.text_frame.text[:50]!r}"
                )


def main() -> int:
    prs = Presentation(str(TEMPLATE))

    slide_title(prs)
    slide_content(prs)
    slide_abstract(prs)
    slide_introduction(prs)
    slide_problem(prs)
    slide_objectives(prs)
    slide_methodology(prs)
    slide_results(prs)
    slide_status(prs)
    slide_publication(prs)

    # Continuations, duplicated from the template slide of the same section.
    method_2 = duplicate(prs, T_METHOD)
    slide_methodology_2(method_2)
    method_3 = duplicate(prs, T_METHOD)
    slide_methodology_3(method_3)
    demo_slides = [prs.slides[T_DEMO]] + [duplicate(prs, T_DEMO) for _ in range(2)]
    for i, (slide, rows) in enumerate(
        zip(demo_slides, SHOT_ROWS, strict=True), start=1
    ):
        slide_demo(slide, rows, i, len(SHOT_ROWS))
    results_2 = duplicate(prs, T_RESULTS)
    slide_results_2(results_2)
    references = duplicate(prs, T_PUB)
    slide_references(references)

    # Template order, with each continuation directly after its own section.
    reorder(prs, [0, 1, 2, 3, 4, 5, 6, 11, 12, 7, 13, 14, 8, 15, 9, 10, 16])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    check(OUT)
    print(f"{OUT.relative_to(ROOT)}  —  {len(prs.slides)} slides")
    return 0


if __name__ == "__main__":
    sys.exit(main())
